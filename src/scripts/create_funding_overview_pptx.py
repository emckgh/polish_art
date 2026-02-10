#!/usr/bin/env python3
"""
Generate a high-level PowerPoint overview of the Polish Looted Art Discovery Engine
for museum directors and government. Professional consulting style (McKinsey/Bain-like):
muted backgrounds, clean typography, optional tasteful artwork image from the DB,
modern headers and footers.

Screenshots: Use scripts/screenshot_for_pptx.py or place PNGs in docs/slideshow/
(screenshot_list.png, screenshot_detail.png).

Requires: pip install python-pptx Pillow
Run: python scripts/create_funding_overview_pptx.py
Output: docs/Polish_Looted_Art_Discovery_Engine_Overview.pptx
"""
import io
import os
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError as e:
    raise SystemExit(f"Install python-pptx: pip install python-pptx (error: {e})")

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = PROJECT_ROOT / "docs" / "Polish_Looted_Art_Discovery_Engine_Overview.pptx"
SLIDESHOW_DIR = PROJECT_ROOT / "docs" / "slideshow"

# Consulting-style palette: muted, professional
OFF_WHITE = RGBColor(248, 248, 250)
BG_LIGHT = RGBColor(252, 252, 253)
CHARCOAL = RGBColor(33, 33, 36)
BODY_GRAY = RGBColor(70, 70, 76)
ACCENT_GRAY = RGBColor(120, 120, 128)
FOOTER_GRAY = RGBColor(100, 100, 108)
LINE_GRAY = RGBColor(200, 200, 204)

# Margins and layout
MARGIN_H = 0.7
MARGIN_TOP = 0.6
MARGIN_BODY_TOP = 1.35
FOOTER_TOP = 6.95
FOOTER_H = 0.35
TITLE_LINE_W = 1.5  # pt
HEADER_LINE_Y = 1.05

SLIDES = [
    {"title": "Polish Looted Art Discovery Engine", "subtitle": "Recovering Cultural Heritage Through Technology", "type": "title"},
    {"title": "Challenge & Mission", "body": [
        "Hundreds of thousands of Polish works looted in WWII still circulate online.",
        "Recovery today is manual and slow.",
        "We use technology to find where looted art appears—auctions, marketplaces, galleries—and surface leads for experts and restitution.",
    ], "type": "content"},
    {"title": "The Platform", "body": [
        "Reference database from official sources (e.g. Division for Looted Art).",
        "Reverse image search at scale; prioritizes auction and marketplace hits.",
        "Web UI for curators to browse, search, and review leads.",
    ], "type": "content_with_images", "images": ["screenshot_list.png", "screenshot_detail.png"], "captions": ["Artwork list", "Artwork detail & Vision results"]},
    {"title": "Impact & Next Steps", "body": [
        "Prioritized leads, low cost, scalable to tens of thousands of artworks.",
        "Pilot with ministry or museum partners; secure funding for API and capacity.",
        "Extend: public reporting, official DB integration, stronger restitution workflows.",
    ], "type": "content"},
    {"title": "Thank You", "body": [
        "We seek partners in museums, government, and cultural institutions to scale impact and accelerate restitution.",
    ], "type": "content"},
]


def _get_one_artwork_image_path() -> Path | None:
    """Try to load one artwork image from DB; return path to a muted, washed-out temp image or None."""
    try:
        from PIL import Image
    except ImportError:
        return None
    db_url = os.environ.get("DATABASE_URL", "").strip() or "sqlite:///data/artworks.db"
    if not db_url.startswith("sqlite"):
        return None
    # Resolve path: sqlite:///data/artworks.db -> PROJECT_ROOT/data/artworks.db
    db_path = db_url.replace("sqlite:///", "").lstrip("/")
    if not os.path.isabs(db_path):
        db_path = str(PROJECT_ROOT / db_path.replace("/", os.sep))
    db_path = str(Path(db_path))
    if not Path(db_path).is_file():
        return None
    try:
        import sqlite3
        conn = sqlite3.connect(db_path)
        row = conn.execute(
            "SELECT image_data, image_mime_type FROM artworks WHERE image_data IS NOT NULL AND length(image_data) > 1000 LIMIT 1"
        ).fetchone()
        conn.close()
        if not row:
            return None
        blob, mime = row
        if not blob:
            return None
        img = Image.open(io.BytesIO(blob)).convert("RGB")
        w, h = img.size
        if w < 200 or h < 200:
            return None
        # Resize to slide-ish size and create very muted version (blend with white)
        img = img.resize((min(1200, w), min(1200, int(h * 1200 / w))), Image.Resampling.LANCZOS)
        white = Image.new("RGB", img.size, (255, 255, 255))
        muted = Image.blend(img, white, alpha=0.88)
        out = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        muted.save(out.name, "PNG")
        out.close()
        return Path(out.name)
    except Exception:
        return None


def _add_muted_background(slide, prs, temp_bg_path: Path | None) -> None:
    """Add a tasteful muted background.

    If a washed-out DB artwork image is available, use it as a subtle full-bleed background.
    Otherwise use a clean, solid light background.
    """
    spTree = slide.shapes._spTree
    if temp_bg_path and temp_bg_path.exists():
        try:
            pic = slide.shapes.add_picture(
                str(temp_bg_path), Inches(0), Inches(0),
                width=prs.slide_width, height=prs.slide_height
            )
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
            return
        except Exception:
            # Fall back to solid background below
            pass

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.fill.background()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_footer(slide, prs, slide_num: int, total: int) -> None:
    """Thin line and footer text in consulting style."""
    # Line
    line = slide.shapes.add_shape(1, Inches(MARGIN_H), Inches(FOOTER_TOP), Inches(10 - 2 * MARGIN_H), Inches(0.003))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_GRAY
    line.line.fill.background()
    # Left: project name
    left = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(FOOTER_TOP + 0.06), Inches(5), Inches(FOOTER_H))
    left.text_frame.text = "Polish Looted Art Discovery Engine"
    left.text_frame.paragraphs[0].font.size = Pt(9)
    left.text_frame.paragraphs[0].font.color.rgb = FOOTER_GRAY
    # Right: slide number
    right = slide.shapes.add_textbox(Inches(10 - MARGIN_H - 1.2), Inches(FOOTER_TOP + 0.06), Inches(1.2), Inches(FOOTER_H))
    right.text_frame.text = f"{slide_num}  |  {total}"
    right.text_frame.paragraphs[0].font.size = Pt(9)
    right.text_frame.paragraphs[0].font.color.rgb = FOOTER_GRAY
    right.text_frame.paragraphs[0].alignment = 2  # right


def _add_header_line(slide, prs) -> None:
    """Thin horizontal accent line under title area."""
    line = slide.shapes.add_shape(1, Inches(MARGIN_H), Inches(HEADER_LINE_Y), Inches(10 - 2 * MARGIN_H), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_GRAY
    line.line.fill.background()


def add_title_slide(prs, data, slide_num, total, temp_bg_path) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_muted_background(slide, prs, temp_bg_path)
    # Title: left-aligned, charcoal
    title_box = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(2.4), Inches(9), Inches(1.1))
    tf = title_box.text_frame
    tf.text = data["title"]
    p = tf.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(3.55), Inches(9), Inches(0.7))
    sub.text_frame.text = data["subtitle"]
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = BODY_GRAY
    _add_footer(slide, prs, slide_num, total)


def add_content_slide(prs, data, slide_num, total, temp_bg_path) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_muted_background(slide, prs, temp_bg_path)
    # Title: left-aligned, charcoal, no big bar
    title_box = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(MARGIN_TOP), Inches(9), Inches(0.55))
    title_box.text_frame.text = data["title"]
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CHARCOAL
    _add_header_line(slide, prs)
    # Body
    body = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(MARGIN_BODY_TOP), Inches(9), Inches(5.2))
    tf = body.text_frame
    tf.text = data["body"][0]
    for line in data["body"][1:]:
        tf.add_paragraph().text = line
    for p in tf.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_GRAY
        p.space_after = Pt(8)
    _add_footer(slide, prs, slide_num, total)


def add_content_with_images_slide(prs, data, slide_num, total, temp_bg_path) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_muted_background(slide, prs, temp_bg_path)
    # Title + line
    title_box = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(MARGIN_TOP), Inches(9), Inches(0.5))
    title_box.text_frame.text = data["title"]
    title_box.text_frame.paragraphs[0].font.size = Pt(22)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = CHARCOAL
    _add_header_line(slide, prs)
    # Body left
    body = slide.shapes.add_textbox(Inches(MARGIN_H), Inches(MARGIN_BODY_TOP), Inches(4.0), Inches(2.0))
    tf = body.text_frame
    tf.text = data["body"][0]
    for line in data["body"][1:]:
        tf.add_paragraph().text = line
    for p in tf.paragraphs:
        p.font.size = Pt(13)
        p.font.color.rgb = BODY_GRAY
        p.space_after = Pt(6)
    # Screenshots right
    SLIDESHOW_DIR.mkdir(parents=True, exist_ok=True)
    img_paths = data.get("images", [])
    captions = data.get("captions", [])
    y = Inches(MARGIN_BODY_TOP)
    img_h = Inches(2.1)
    img_w = Inches(4.9)
    x_img = Inches(4.85)
    for i, name in enumerate(img_paths):
        path = SLIDESHOW_DIR / name
        if path.exists():
            try:
                slide.shapes.add_picture(str(path), x_img, y, width=img_w, height=img_h)
            except Exception:
                _add_placeholder(slide, x_img, y, img_w, img_h, captions[i] if i < len(captions) else name)
        else:
            _add_placeholder(slide, x_img, y, img_w, img_h, captions[i] if i < len(captions) else name)
        y += img_h + Inches(0.12)
    _add_footer(slide, prs, slide_num, total)


def _add_placeholder(slide, left, top, width, height, caption) -> None:
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 245, 247)
    box.line.color.rgb = LINE_GRAY
    tx = slide.shapes.add_textbox(left, top + height / 2 - Inches(0.15), width, Inches(0.3))
    tx.text_frame.text = f"[Screenshot: {caption}]"
    tx.text_frame.paragraphs[0].font.size = Pt(9)
    tx.text_frame.paragraphs[0].font.color.rgb = ACCENT_GRAY
    tx.text_frame.paragraphs[0].alignment = 1


def main() -> None:
    temp_bg_path = _get_one_artwork_image_path()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    total = len(SLIDES)

    for idx, data in enumerate(SLIDES, start=1):
        if data["type"] == "title":
            add_title_slide(prs, data, idx, total, temp_bg_path)
        elif data["type"] == "content_with_images":
            add_content_with_images_slide(prs, data, idx, total, temp_bg_path)
        else:
            add_content_slide(prs, data, idx, total, temp_bg_path)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_PATH
    try:
        prs.save(str(out_path))
    except PermissionError:
        # If the target file is open/locked (common on Windows), fall back to a new filename.
        base_new = OUTPUT_PATH.parent / (OUTPUT_PATH.stem + "_new.pptx")
        try:
            out_path = base_new
            prs.save(str(out_path))
            print("(Original file is open; saved as _new.pptx instead.)")
        except PermissionError:
            import datetime as _dt
            stamp = _dt.datetime.now().strftime("%Y%m%d_%H%M%S")
            out_path = OUTPUT_PATH.parent / f"{OUTPUT_PATH.stem}_new_{stamp}.pptx"
            prs.save(str(out_path))
            print("(Original and _new files are open; saved with timestamp instead.)")
    print(f"Saved: {out_path}")
    if temp_bg_path and temp_bg_path.exists():
        try:
            temp_bg_path.unlink(missing_ok=True)
        except Exception:
            pass


if __name__ == "__main__":
    main()
